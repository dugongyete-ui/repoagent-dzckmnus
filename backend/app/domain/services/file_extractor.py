"""
Server-side file text extractor.

Handles .pptx, .pdf, .docx, .xlsx, .csv and plain text files
by extracting their text content in Python — no shell commands,
no AI guesswork, no sandbox needed.
"""
import io
import logging
from typing import Optional

logger = logging.getLogger(__name__)

# Supported binary formats and their MIME types
EXTRACTABLE_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # pptx
    "application/vnd.ms-powerpoint",  # ppt (limited)
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",  # docx
    "application/msword",  # doc
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",  # xlsx
    "application/vnd.ms-excel",  # xls
    "text/csv",
    "text/plain",
    "text/markdown",
    "text/x-python",
    "application/json",
    "application/xml",
    "text/xml",
    "text/html",
}

EXTRACTABLE_EXTENSIONS = {
    ".pptx", ".ppt", ".pdf", ".docx", ".doc",
    ".xlsx", ".xls", ".csv", ".txt", ".md",
    ".py", ".js", ".ts", ".json", ".xml", ".html", ".htm",
}


def is_extractable(filename: str, content_type: Optional[str] = None) -> bool:
    """Return True if this file can be text-extracted server-side."""
    if filename:
        ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
        if ext in EXTRACTABLE_EXTENSIONS:
            return True
    if content_type:
        base = content_type.lower().split(";")[0].strip()
        if base in EXTRACTABLE_MIME_TYPES or base.startswith("text/"):
            return True
    return False


def extract_text(data: bytes, filename: str, content_type: Optional[str] = None) -> str:
    """
    Extract plain text from file bytes.

    Returns the extracted text string, or raises an exception if extraction fails.
    Caller should catch and log — never crash the agent over a file.
    """
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""
    mime = (content_type or "").lower().split(";")[0].strip()

    # ── PowerPoint (.pptx) ────────────────────────────────────────────────
    if ext in (".pptx", ".ppt") or "presentationml" in mime:
        return _extract_pptx(data, filename)

    # ── PDF ───────────────────────────────────────────────────────────────
    if ext == ".pdf" or mime == "application/pdf":
        return _extract_pdf(data, filename)

    # ── Word (.docx) ──────────────────────────────────────────────────────
    if ext in (".docx", ".doc") or "wordprocessingml" in mime or mime == "application/msword":
        return _extract_docx(data, filename)

    # ── Excel (.xlsx / .xls) ─────────────────────────────────────────────
    if ext in (".xlsx", ".xls") or "spreadsheetml" in mime or mime in ("application/vnd.ms-excel",):
        return _extract_xlsx(data, filename)

    # ── CSV ───────────────────────────────────────────────────────────────
    if ext == ".csv" or mime == "text/csv":
        return _extract_csv(data)

    # ── Plain text / code / markdown ─────────────────────────────────────
    if ext in (".txt", ".md", ".py", ".js", ".ts", ".json", ".xml", ".html", ".htm") \
            or mime.startswith("text/") or mime in ("application/json", "application/xml"):
        return _extract_text(data)

    raise ValueError(f"Unsupported file type: ext={ext!r} mime={mime!r}")


# ── Per-format extractors ─────────────────────────────────────────────────────

def _extract_pptx(data: bytes, filename: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        lines = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                lines.append(f"[Slide {slide_num}]")
                lines.extend(slide_texts)
        text = "\n".join(lines)
        logger.info(f"PPTX extracted: {filename} → {len(lines)} blocks, {len(text)} chars")
        return text
    except Exception as e:
        logger.error(f"PPTX extraction failed for {filename}: {e}")
        raise


def _extract_pdf(data: bytes, filename: str) -> str:
    try:
        import pdfplumber
        pages = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    pages.append(f"[Page {i}]\n{text.strip()}")
        text = "\n\n".join(pages)
        logger.info(f"PDF extracted: {filename} → {len(pages)} pages, {len(text)} chars")
        return text
    except Exception as e:
        logger.error(f"PDF extraction failed for {filename}: {e}")
        raise


def _extract_docx(data: bytes, filename: str) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)
        logger.info(f"DOCX extracted: {filename} → {len(paragraphs)} paragraphs, {len(text)} chars")
        return text
    except Exception as e:
        logger.error(f"DOCX extraction failed for {filename}: {e}")
        raise


def _extract_xlsx(data: bytes, filename: str) -> str:
    try:
        import pandas as pd
        sheets = pd.read_excel(io.BytesIO(data), sheet_name=None)
        parts = []
        for sheet_name, df in sheets.items():
            parts.append(f"[Sheet: {sheet_name}]\n{df.to_string()}")
        text = "\n\n".join(parts)
        logger.info(f"XLSX extracted: {filename} → {len(sheets)} sheets, {len(text)} chars")
        return text
    except Exception as e:
        logger.error(f"XLSX extraction failed for {filename}: {e}")
        raise


def _extract_csv(data: bytes) -> str:
    try:
        import pandas as pd
        df = pd.read_csv(io.BytesIO(data))
        return df.to_string()
    except Exception:
        return data.decode("utf-8", errors="replace")


def _extract_text(data: bytes) -> str:
    return data.decode("utf-8", errors="replace")
